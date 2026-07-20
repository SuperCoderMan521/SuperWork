#!/usr/bin/env python3
"""Generate a professional PPT about New Energy Vehicle (NEV) trends."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color Palette ──
BG_DARK    = RGBColor(0x1A, 0x2A, 0x3A)   # deep navy
BG_CARD    = RGBColor(0x22, 0x3A, 0x4E)   # slightly lighter
ACCENT     = RGBColor(0x00, 0xD4, 0xAA)   # teal / mint green
ACCENT2    = RGBColor(0x4A, 0xC8, 0xFF)   # sky blue
ACCENT_ORANGE = RGBColor(0xFF, 0x9F, 0x43)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xDD, 0xEE)
MED_GRAY   = RGBColor(0x88, 0xAA, 0xBB)
DARK_TEXT   = RGBColor(0x0A, 0x1A, 0x2A)
GREEN_BG    = RGBColor(0xE8, 0xF5, 0xE9)

# ── Helper Functions ──

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_with_text(slide, left, top, width, height, text, font_size=18,
                        color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                        font_name="Microsoft YaHei", fill_color=None,
                        shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    # vertical centering
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=WHITE, font_name="Microsoft YaHei", line_spacing=1.5):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(6)
        p.level = 0
        # bullet character
        pPr = p._pPr
        if pPr is None:
            from lxml import etree
            pPr = etree.SubElement(p._p, '{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
        # We'll just use a simple dash prefix
        p.text = f"▸  {item}"
    return txBox

def add_card(slide, left, top, width, height, title, body, title_size=18,
             body_size=14, title_color=ACCENT, body_color=LIGHT_GRAY,
             card_color=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = card_color
    shape.line.fill.background()
    # title
    add_textbox(slide, left + Inches(0.3), top + Inches(0.2), width - Inches(0.6), Inches(0.5),
                title, font_size=title_size, color=title_color, bold=True)
    # body
    add_textbox(slide, left + Inches(0.3), top + Inches(0.7), width - Inches(0.6), height - Inches(0.9),
                body, font_size=body_size, color=body_color)

def add_accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.5), color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_page_number(slide, num, total):
    add_textbox(slide, Inches(12.5), Inches(7.0), Inches(0.7), Inches(0.4),
                f"{num}/{total}", font_size=10, color=MED_GRAY,
                alignment=PP_ALIGN.RIGHT)

def add_section_title(slide, title, subtitle=None):
    add_accent_bar(slide, Inches(0.6), Inches(0.5), height=Inches(0.6))
    add_textbox(slide, Inches(0.9), Inches(0.4), Inches(10), Inches(0.7),
                title, font_size=28, color=WHITE, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.9), Inches(1.1), Inches(10), Inches(0.5),
                    subtitle, font_size=14, color=MED_GRAY)

TOTAL_SLIDES = 12

# ════════════════════════════════════════════
# Slide 1 — Title
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG_DARK)

# large accent shape
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                 Inches(0.15), Inches(7.5))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

# title
add_textbox(slide, Inches(1.2), Inches(1.8), Inches(11), Inches(1.2),
            "新能源汽车发展趋势", font_size=48, color=WHITE, bold=True,
            alignment=PP_ALIGN.LEFT)

# subtitle
add_textbox(slide, Inches(1.2), Inches(3.2), Inches(11), Inches(0.8),
            "New Energy Vehicle — 2025-2030 全球市场与技术展望",
            font_size=24, color=ACCENT, alignment=PP_ALIGN.LEFT)

# decorative line
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(4.2),
                                 Inches(3), Inches(0.06))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

# info
add_textbox(slide, Inches(1.2), Inches(4.6), Inches(6), Inches(0.5),
            "2025年7月  |  行业分析报告", font_size=16, color=MED_GRAY)

# small decorative elements on the right
for i, c in enumerate([ACCENT, ACCENT2, ACCENT_ORANGE]):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(11.5 + i * 0.5), Inches(6.0),
                                     Inches(0.25), Inches(0.25))
    shape.fill.solid()
    shape.fill.fore_color.rgb = c
    shape.line.fill.background()

add_page_number(slide, 1, TOTAL_SLIDES)

# ════════════════════════════════════════════
# Slide 2 — 目录
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_title(slide, "目  录", "CONTENTS")

toc_items = [
    ("01", "全球新能源汽车市场概览"),
    ("02", "核心驱动因素分析"),
    ("03", "技术路线对比：BEV / PHEV / FCEV"),
    ("04", "电池技术发展"),
    ("05", "充电基础设施"),
    ("06", "主要市场格局：中国 · 欧洲 · 美国"),
    ("07", "竞争格局与主要玩家"),
    ("08", "智能化 + 电动化融合趋势"),
    ("09", "未来展望与预测"),
    ("10", "总结与建议"),
]

for i, (num, title) in enumerate(toc_items):
    y = Inches(1.8) + Inches(i * 0.5)
    add_textbox(slide, Inches(1.2), y, Inches(0.7), Inches(0.5),
                num, font_size=20, color=ACCENT, bold=True,
                alignment=PP_ALIGN.RIGHT)
    add_textbox(slide, Inches(2.0), y, Inches(8), Inches(0.5),
                title, font_size=18, color=WHITE)

add_page_number(slide, 2, TOTAL_SLIDES)

# ════════════════════════════════════════════
# Slide 3 — 全球市场概览
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_title(slide, "全球新能源汽车市场概览", "Global NEV Market Overview")

cards_data = [
    ("2024年全球销量", "超 1,750 万辆\n同比增长约 25%\n渗透率达 22%", ACCENT),
    ("中国市场份额", "占比全球 60%+\n销量超 1,100 万辆\n渗透率突破 45%", ACCENT2),
    ("欧洲市场", "销量约 300 万辆\n渗透率约 24%\n增速趋于平稳", ACCENT_ORANGE),
    ("美国市场", "销量约 160 万辆\n渗透率约 10%\n政策驱动增长", RGBColor(0xE8, 0x6C, 0xA8)),
]

card_w = Inches(2.8)
card_h = Inches(3.0)
gap = Inches(0.3)
start_x = Inches(0.7)
start_y = Inches(2.0)

for i, (title, body, clr) in enumerate(cards_data):
    x = start_x + i * (card_w + gap)
    add_card(slide, x, start_y, card_w, card_h, title, body,
             title_color=clr, body_color=LIGHT_GRAY)
    # small accent line on top of card
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, start_y, card_w, Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = clr
    shape.line.fill.background()

# key insight
add_textbox(slide, Inches(0.7), Inches(5.5), Inches(12), Inches(0.8),
            "核心洞察：全球新能源汽车已从'早期采用者'阶段进入'早期大众'阶段，"
            "中国市场引领全球，2025 年渗透率有望突破 50%。",
            font_size=15, color=ACCENT)

add_page_number(slide, 3, TOTAL_SLIDES)

# ════════════════════════════════════════════
# Slide 4 — 核心驱动因素
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_title(slide, "核心驱动因素分析", "Key Drivers")

drivers = [
    ("🌍  碳中和目标", "全球 130+ 国家提出碳中和目标，\n交通领域电气化是关键路径。\n欧盟 2035 年禁售燃油车已立法。"),
    ("🔋  电池成本下降", "电池组成本从 2013 年 $600/kWh\n降至 2024 年约 $115/kWh。\n未来有望突破 $70/kWh 门槛。"),
    ("🏛️  政策推动", "购车补贴、绿牌路权、\n双积分政策、碳排放法规\n等多重政策合力驱动。"),
    ("📱  用户接受度提升", "续航焦虑大幅缓解（600km+），\n使用成本低至燃油车 1/5，\n智能化体验吸引年轻用户。"),
]

for i, (title, body) in enumerate(drivers):
    x = Inches(0.7) + (i % 2) * Inches(6.3)
    y = Inches(2.0) + (i // 2) * Inches(2.5)
    add_card(slide, x, y, Inches(5.8), Inches(2.1), title, body,
             title_color=ACCENT, body_color=LIGHT_GRAY,
             card_color=BG_CARD)

add_page_number(slide, 4, TOTAL_SLIDES)

# ════════════════════════════════════════════
# Slide 5 — 技术路线对比
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_title(slide, "技术路线对比", "BEV / PHEV / FCEV")

techs = [
    ("BEV 纯电动", "🔋", "当前主流，技术最成熟\n续航 400-800km\n效率最高，TCO 最低\n代表：Tesla Model 3/Y,\n比亚迪秦/汉/海豹", ACCENT),
    ("PHEV 插电混动", "⚡", "过渡方案，仍占较大份额\n纯电续航 50-200km\n兼顾长途需求\n代表：比亚迪 DM-i,\n理想 L 系列", ACCENT2),
    ("FCEV 氢燃料", "💧", "远期潜力技术\n续航 500-700km\n加氢 3-5 分钟\n基础设施是最大瓶颈\n代表：Toyota Mirai", ACCENT_ORANGE),
]

for i, (title, icon, body, clr) in enumerate(techs):
    x = Inches(0.7) + i * Inches(4.2)
    add_card(slide, x, Inches(2.0), Inches(3.8), Inches(4.5), f"{icon}  {title}", body,
             title_color=clr, body_color=LIGHT_GRAY,
             card_color=BG_CARD)
    # header accent
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.0), Inches(3.8), Inches(0.06))
    shape.fill.solid()
    shape.fill.fore_color.rgb = clr
    shape.line.fill.background()

add_page_number(slide, 5, TOTAL_SLIDES)

# ════════════════════════════════════════════
# Slide 6 — 电池技术
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_title(slide, "电池技术发展", "Battery Technology Evolution")

# left - timeline
add_textbox(slide, Inches(0.7), Inches(2.0), Inches(5.5), Inches(0.5),
            "技术演进路线", font_size=20, color=ACCENT, bold=True)

timeline = [
    ("2020-2023", "磷酸铁锂（LFP）崛起 — 成本低、安全性高"),
    ("2023-2025", "三元锂（NCM）高镍化 — 能量密度 300Wh/kg+"),
    ("2024-2026", "固态电池量产起步 — 丰田/宁德时代/QuantumScape"),
    ("2026-2028", "半固态 → 全固态过渡 — 能量密度 500Wh/kg 目标"),
    ("2028-2030", "钠离子电池 / 锂硫电池商业化 — 进一步降本"),
]

for i, (year, desc) in enumerate(timeline):
    y = Inches(2.7) + i * Inches(0.75)
    # dot
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.85), y + Inches(0.1),
                                     Inches(0.16), Inches(0.16))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT if i < 3 else ACCENT2
    shape.line.fill.background()
    # line
    if i < len(timeline) - 1:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.92), y + Inches(0.26),
                                         Inches(0.03), Inches(0.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = MED_GRAY
        shape.line.fill.background()
    # year
    add_textbox(slide, Inches(1.4), y, Inches(1.8), Inches(0.35),
                year, font_size=14, color=ACCENT if i < 3 else ACCENT2, bold=True)
    # desc
    add_textbox(slide, Inches(3.2), y, Inches(4), Inches(0.35),
                desc, font_size=13, color=LIGHT_GRAY)

# right - key metrics
add_textbox(slide, Inches(8.0), Inches(2.0), Inches(4.5), Inches(0.5),
            "关键指标", font_size=20, color=ACCENT_ORANGE, bold=True)

metrics = [
    ("电池成本 ↓", "$115/kWh (2024) → $70/kWh (2028 目标)"),
    ("能量密度 ↑", "250 → 500 Wh/kg (2024→2030)"),
    ("充电速度 ↑", "30min (10-80%) → 10min (800V 架构)"),
    ("循环寿命 ↑", "2,000 次 → 5,000 次以上"),
]

for i, (title, val) in enumerate(metrics):
    y = Inches(2.7) + i * Inches(1.0)
    add_card(slide, Inches(8.0), y, Inches(4.5), Inches(0.8), title, val,
             title_color=ACCENT_ORANGE, body_color=LIGHT_GRAY,
             card_color=BG_CARD)

add_page_number(slide, 6, TOTAL_SLIDES)

# ════════════════════════════════════════════
# Slide 7 — 充电基础设施
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_title(slide, "充电基础设施", "Charging Infrastructure")

# stats row
stats = [
    ("全球公共充电桩", "450 万+\n(2024)", ACCENT),
    ("中国占全球", "70%+\n约 320 万根", ACCENT2),
    ("快充比例", "约 35%\n并持续提升", ACCENT_ORANGE),
    ("车桩比（全球）", "约 6.5:1\n目标 3:1", RGBColor(0xE8, 0x6C, 0xA8)),
]

for i, (title, val, clr) in enumerate(stats):
    x = Inches(0.7) + i * Inches(3.2)
    add_card(slide, x, Inches(2.0), Inches(2.9), Inches(1.8), title, val,
             title_color=clr, body_color=WHITE, card_color=BG_CARD)

# key trends
add_textbox(slide, Inches(0.7), Inches(4.3), Inches(12), Inches(0.5),
            "关键趋势", font_size=20, color=ACCENT, bold=True)

trends = [
    "800V 高压快充成为新车型标配，充电 10 分钟续航 400km+",
    "超充网络竞争白热化：Tesla Supercharger → 开放 / 华为液冷超充 / 蔚来换电",
    "V2G (Vehicle-to-Grid) 技术兴起，电动车参与电网调度",
    "无线充电技术走向量产，自动泊车 + 无线充体验",
]

add_bullet_list(slide, Inches(0.7), Inches(4.9), Inches(11), Inches(2.2),
                trends, font_size=15, color=LIGHT_GRAY)

add_page_number(slide, 7, TOTAL_SLIDES)

# ════════════════════════════════════════════
# Slide 8 — 主要市场格局
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_title(slide, "主要市场格局", "China · Europe · United States")

markets = [
    ("🇨🇳  中国", "市场地位：全球领导者",
     [
         "2024 年新能源渗透率突破 45%",
         "比亚迪、吉利、长安等本土品牌主导",
         "价格战激烈，加速行业洗牌",
         "政策支持持续但逐步退坡",
         "出口欧洲/东南亚快速增长",
     ]),
    ("🇪🇺  欧洲", "市场地位：全球第二大市场",
     [
         "2024 年渗透率约 24%，增速放缓",
         "德国取消补贴后短期承压",
         "Stellantis/VW 加速电动化转型",
         "中国品牌（MG、BYD）份额提升",
         "2035 禁燃令奠定长期确定性",
     ]),
    ("🇺🇸  美国", "市场地位：高潜力市场",
     [
         "2024 年渗透率约 10% 较低",
         "IRA 法案推动本土电池供应链",
         "Tesla 主导，传统 OEM 加速追赶",
         "充电基础设施仍为瓶颈",
         "2025-2027 年多款平价车型上市",
     ]),
]

for i, (title, subtitle, items) in enumerate(markets):
    x = Inches(0.5) + i * Inches(4.3)
    # card
    add_card(slide, x, Inches(2.0), Inches(3.9), Inches(4.8), title, "",
             title_color=WHITE, body_color=LIGHT_GRAY, card_color=BG_CARD)
    # subtitle in card
    add_textbox(slide, x + Inches(0.3), Inches(2.6), Inches(3.3), Inches(0.4),
                subtitle, font_size=11, color=MED_GRAY)
    # bullet items
    for j, item in enumerate(items):
        y = Inches(3.1) + j * Inches(0.6)
        add_textbox(slide, x + Inches(0.3), y, Inches(3.3), Inches(0.5),
                    f"•  {item}", font_size=12, color=LIGHT_GRAY)

add_page_number(slide, 8, TOTAL_SLIDES)

# ════════════════════════════════════════════
# Slide 9 — 竞争格局
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_title(slide, "竞争格局与主要玩家", "Competitive Landscape")

players = [
    ("特斯拉 Tesla", "全球 BEV 龙头\nFSD + 超充网络护城河\n2024 全球销量 ~180 万辆\nCybertruck / Model 2 新增长极", ACCENT),
    ("比亚迪 BYD", "中国霸主 + 全球扩张\n垂直整合（电池/芯片/整车）\nDM-i 插混 + 纯电双轮驱动\n2024 全球销量 ~420 万辆（含 PHEV）", ACCENT2),
    ("新势力", "蔚来：换电 + 高端化\n小鹏：智驾技术领先\n理想：增程大空间精准定位\n小米：生态链降维打击", ACCENT_ORANGE),
    ("传统车企转型", "大众：ID 系列 + SSP 平台\n丰田：氢能 + 电动双重路线\n通用：Ultium 平台 + 中美双市场\nStellantis：STLA 平台 + 合作", RGBColor(0xE8, 0x6C, 0xA8)),
]

for i, (title, body, clr) in enumerate(players):
    x = Inches(0.5) + (i % 2) * Inches(6.4)
    y = Inches(2.0) + (i // 2) * Inches(2.5)
    add_card(slide, x, y, Inches(6.0), Inches(2.1), title, body,
             title_color=clr, body_color=LIGHT_GRAY, card_color=BG_CARD)

add_page_number(slide, 9, TOTAL_SLIDES)

# ════════════════════════════════════════════
# Slide 10 — 智能化 + 电动化
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_title(slide, "智能化 + 电动化融合趋势", "Software-Defined Vehicle")

areas = [
    ("智能驾驶 (ADAS/AD)", "L2+ 成标配，L3 开始落地\n城市 NOA 开城竞赛\nTesla FSD / 华为 ADS / XNGP\nBEV + Transformer 端到端方案"),
    ("智能座舱", "AI 大模型上车\n语音助手 + 多模态交互\n车内生态：办公/娱乐/休息\n座舱芯片：高通 8295 / AMD"),
    ("OTA + 软件服务", "整车 OTA 成标配（SOTA/FOTA）\n软件订阅收入成新利润点\nTesla FSD $99/月 模式验证\n自动驾驶数据闭环"),
    ("车联网 (V2X)", "5G + C-V2X 车路协同\n高精地图 + 实时路况\nOTA 地图 + 云端调度\n智慧城市基础设施联动"),
]

for i, (title, body) in enumerate(areas):
    x = Inches(0.5) + (i % 2) * Inches(6.4)
    y = Inches(2.0) + (i // 2) * Inches(2.5)
    add_card(slide, x, y, Inches(6.0), Inches(2.1), title, body,
             title_color=ACCENT2, body_color=LIGHT_GRAY, card_color=BG_CARD)

add_page_number(slide, 10, TOTAL_SLIDES)

# ════════════════════════════════════════════
# Slide 11 — 未来展望
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)
add_section_title(slide, "未来展望与预测", "Outlook & Predictions")

# left - predictions
predictions = [
    ("2025", "全球新能源渗透率突破 25%\n中国渗透率稳定在 50%+ 水平"),
    ("2027", "800V 快充成为 15 万以上车型标配\nL3 自动驾驶在多国合法化"),
    ("2030", "全球新能源车销量占比达 40-50%\n电池成本降至 $70/kWh 以下\n全固态电池开始规模应用"),
    ("2035", "欧盟/多国禁售燃油车生效\n新能源渗透率超 70%\n氢燃料在商用车领域开始普及"),
]

for i, (year, desc) in enumerate(predictions):
    y = Inches(2.0) + i * Inches(1.2)
    # year badge
    add_shape_with_text(slide, Inches(0.7), y, Inches(1.2), Inches(0.45),
                        year, font_size=16, color=BG_DARK, bold=True,
                        fill_color=ACCENT, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2.1), y, Inches(5), Inches(0.9),
                desc, font_size=13, color=LIGHT_GRAY)

# right - key insight
add_card(slide, Inches(8.0), Inches(2.0), Inches(4.5), Inches(4.5),
         "💡  核心观点",
         "未来 5 年是新能源汽车行业的'分水岭'：\n\n"
         "• 技术端：电池成本拐点已至，电动化\n   vs 燃油车平价时代来临\n\n"
         "• 市场端：中国品牌出海加速，全球\n   格局重塑\n\n"
         "• 用户端：智能化体验成为差异化\n   核心竞争力\n\n"
         "• 产业端：价值链重构，软件和服务\n   收入占比持续提升",
         title_color=ACCENT, body_color=LIGHT_GRAY, card_color=BG_CARD)

add_page_number(slide, 11, TOTAL_SLIDES)

# ════════════════════════════════════════════
# Slide 12 — 总结
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

# large "Thank You" text
add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
            "感谢聆听", font_size=56, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.8),
            "Thank You", font_size=32, color=ACCENT,
            alignment=PP_ALIGN.CENTER)

# decorative line
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.2),
                                 Inches(2.3), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

# key takeaway
add_textbox(slide, Inches(1.5), Inches(4.6), Inches(10), Inches(0.8),
            "新能源汽车行业正处于从'政策驱动'到'产品驱动'的关键转折期\n"
            "电动化是上半场，智能化是下半场",
            font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# summary points
points = [
    "全球新能源渗透率持续攀升，2025 年有望突破 25%",
    "中国引领全球，本土品牌加速出海",
    "电池成本下降 + 智能化体验 = 核心增长双引擎",
]

for i, pt in enumerate(points):
    y = Inches(5.5) + i * Inches(0.5)
    add_textbox(slide, Inches(2), y, Inches(9), Inches(0.5),
                f"✦  {pt}", font_size=14, color=MED_GRAY,
                alignment=PP_ALIGN.CENTER)

add_page_number(slide, 12, TOTAL_SLIDES)

# ── Save ──
output_path = os.path.expanduser("~/Desktop/新能源汽车发展趋势.pptx")
prs.save(output_path)
print(f"[OK] PPT saved to: {output_path}")
